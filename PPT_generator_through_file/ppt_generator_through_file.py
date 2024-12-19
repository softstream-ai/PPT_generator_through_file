import openai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import uvicorn
import docx
from PyPDF2 import PdfReader
from fastapi import FastAPI, Form, HTTPException, UploadFile
from fastapi.responses import FileResponse
from pptx.enum.text import PP_ALIGN
from fastapi.middleware.cors import CORSMiddleware


# Set up OpenAI API key
 # Replace with your OpenAI API key



def extract_text_from_word(file):
    """Extract text from Word documents."""
    doc = docx.Document(file)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])


def extract_text_from_pdf(file):
    """Extract text from PDF documents."""
    reader = PdfReader(file)
    return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])


def generate_title(file_content):
    """Generate a title for the presentation using OpenAI."""
    try:
        messages = [
            {"role": "system", "content": "You are a helpful assistant generating presentation titles."},
            {"role": "user", "content": f"Generate a concise, descriptive title for a presentation based on this content:\n{file_content}"}
        ]
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=messages,
            max_tokens=50,
            temperature=0.7
        )
        return response['choices'][0]['message']['content'].strip()
    except Exception as e:
        return "Generated Presentation"  # Fallback title in case of an error


def generate_subtitle(file_content):
    """Generate a subtitle for the title slide using OpenAI."""
    try:
        messages = [
            {"role": "system", "content": "You are a helpful assistant generating subtitles for presentation slides."},
            {"role": "user", "content": f"Generate a subtitle for a presentation based on this content:\n{file_content}"}
        ]
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=messages,
            max_tokens=50,
            temperature=0.7
        )
        return response['choices'][0]['message']['content'].strip()
    except Exception as e:
        return "Subtitle could not be generated."


def generate_presentation_content(file_content):
    """Generate content for the presentation using OpenAI."""
    try:
        messages = [
            {"role": "system", "content": "You are a helpful assistant generating presentation outlines."},
            {"role": "user", "content": f"Generate a detailed presentation outline from:\n{file_content}. Generate a title, section titles, and for each section, provide 3 numbered points followed by their definitions (explain) in lines as bullet points. Ensure that each point is numbered (1, 2, 3) and the explanations of the points are in bullet format. Create at least 15 slide points."},
    # Q&A slide
            {"role": "user", "content": f"Add a Future of content slide to encourage audience interaction."},
    # Thank you slide
           # {"role": "user", "content": f"Add a 'Thank You' slide to conclude the presentation."}

        ]
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=messages,
            max_tokens=5000,
            temperature=0.7
        )
        return response['choices'][0]['message']['content'].strip()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"OpenAI Error: {str(e)}")


def apply_template(template_choice):
    templates = {
        "Business1": "Business1.pptx",
        "Business2": "Business2.pptx",
        "Business3": "Business3.pptx",
        "Business4": "Business4.pptx",
        "Business5": "Business5.pptx",
        "Student1":"Student1.pptx",
        "Student2":"Student2.pptx",
        "Student3":"Student3.pptx",
        "Student4":"Student4.pptx",
        "Student5":"Student5.pptx",
        "Corporate1":"Corporate1.pptx",
        "Corporate2":"Corporate2.pptx",
        "Corporate3":"Corporate3.pptx",
        "Corporate4":"Corporate4.pptx",
        "Corporate5":"Corporate5.pptx",
        "Creative1":"Creative1.pptx",
        "Creative2":"Creative2.pptx",
        "Creative3":"Creative3.pptx",
        "Creative4":"Creative4.pptx",
        "Creative5":"Creative5.pptx",
    }
    template_path = templates.get(template_choice)
    if template_path and os.path.exists(template_path):
        return Presentation(template_path)  # Load the chosen template
    return Presentation()  # Fallback to a blank presentation



# Function to add a slide with customized text placement
def add_slide(prs, title_text, content_text, template_choice, image_path=None):
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title shape
    title_shape = slide.shapes.title
    title_shape.text = title_text
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)

    # Change the title color based on the template choice
    if template_choice == "Business4":
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 0)  # Yellow
    elif template_choice in ["Creative3","Creative6"]:
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(34, 139, 34)  # Green
    elif template_choice in ["Creative2", "Student5", "Business1", "Business5","Corporate1","Corporate2","Creative5","Corporate3","Corporate4"]:
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White
    else:
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Default navy


    if template_choice in ["Corporate2","Creative5"]:
        text_left = Inches(6)
        text_width = Inches(7)
        content_box = slide.shapes.add_textbox(text_left, Inches(1.5), text_width, Inches(5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

    else:
    # Default Content text box placement and size
        text_left = Inches(1)
        text_width = Inches(11)
        content_box = slide.shapes.add_textbox(text_left, Inches(2.5), text_width, Inches(5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True



    # Add content with bullet points and formatting
    bullet_points = content_text.split('\n')
    for point in bullet_points:
        content_paragraph = content_frame.add_paragraph()
        content_paragraph.text = point.strip()

        # Bold numbered bullet points
        if point.strip().split('.')[0].isdigit():
            content_paragraph.font.bold = True

        content_paragraph.font.size = Pt(18)
        content_paragraph.space_after = Pt(10)
        content_paragraph.alignment = PP_ALIGN.LEFT

    # Add image if provided
    #if image_path and os.path.exists(image_path):
        #img_width = Inches(5.30)
        #img_height = Inches(2.84)

        #slide_width = prs.slide_width
        #slide_height = prs.slide_height

        # Center the image vertically and horizontally
        #top = (slide_height - img_height) / 2 + Inches(1.5)
        #left = (slide_width - img_width) / 2
        #slide.shapes.add_picture(image_path, left, top, width=img_width, height=img_height)


            # Add image if provided
    if image_path and os.path.exists(image_path):
        img_width = Inches(5.30)
        img_height = Inches(2.84)

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Check for specific templates
        if template_choice in ["Corporate2", "Creative5"]:
            # Position the image on the right side and slightly upward
            top = (slide_height - img_height) / 2 + Inches(0.5)  # Move upward by 0.5 inches
            left = slide_width - img_width - Inches(2)           # Align to the right with a 1-inch margin
        else:
            # Default centered placement
            top = (slide_height - img_height) / 2 + Inches(1.5)
            left = (slide_width - img_width) / 2

        slide.shapes.add_picture(image_path, left, top, width=img_width, height=img_height)



    

# Function to create the presentation
def create_presentation(file_content, slide_count, template_choice):
    prs = apply_template(template_choice)

    # Generate the title and subtitle
    topic = generate_title(file_content)  # Generate title

    # Create the title slide
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    title_slide = prs.slides.add_slide(title_slide_layout)

    # Set the title in the title box
    title_shape = title_slide.shapes.title
    title_shape.text = topic
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)

    # Generate subtitle using ChatGPT and set in the subtitle box
    # subtitle = generate_subtitle(file_content)
    # subtitle_shape = title_slide.shapes.placeholders[0]  # Placeholder for subtitle
    # subtitle_shape.text = subtitle
    # subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    # subtitle_shape.text_frame.paragraphs[0].font.italic = True
    # content_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(150, 0, 150)   # Light purple color

    # Generate content for the presentation
    content = generate_presentation_content(file_content)
    sections = content.split("\n\n")[:slide_count]

    # Add the Table of Contents slide (working as per your code)
    toc_slide_layout = prs.slide_layouts[5]  # Blank layout
    toc_slide = prs.slides.add_slide(toc_slide_layout)

    # Add TOC title
    toc_title_box = toc_slide.shapes.title
    toc_title_box.text = "Table of Contents"
    toc_title_box.text_frame.paragraphs[0].font.size = Pt(32)
    toc_title_box.text_frame.paragraphs[0].font.bold = True


    if template_choice in ["Corporate2","Creative5"]:
        text_left = Inches(6)  # Left margin
        text_width = Inches(7)  # Width matching content slides
        toc_content_box = toc_slide.shapes.add_textbox(text_left, Inches(1.5), text_width, Inches(5))
        toc_content_frame = toc_content_box.text_frame
        toc_content_frame.word_wrap = True

    else:
    # Add TOC content box
        text_left = Inches(1)  # Left margin
        text_width = Inches(11)  # Width matching content slides
        toc_content_box = toc_slide.shapes.add_textbox(text_left, Inches(2.5), text_width, Inches(5))
        toc_content_frame = toc_content_box.text_frame
        toc_content_frame.word_wrap = True

    # Add sections to the TOC
    for i, section in enumerate(sections):
        lines = section.strip().split("\n")
        slide_title = lines[0]  # Title of the slide

        # Add TOC entry
        toc_paragraph = toc_content_frame.add_paragraph()
        toc_paragraph.text = f"{slide_title}"
        toc_paragraph.font.size = Pt(18)
        toc_paragraph.space_after = Pt(10)
        toc_paragraph.alignment = PP_ALIGN.LEFT

    # Add slides for each section
    for section in sections:
        lines = section.strip().split("\n")
        slide_title = lines[0]
        slide_content = "\n".join(lines[1:])

        # Add slide with title and content
        add_slide(prs, slide_title, slide_content, template_choice)

    if template_choice in ["Corporate2", "Creative6"]:
        # Add Q&A Slide with specific image for these templates
        add_slide(prs, "Q&A", "Questions & Answers", template_choice, image_path="QandA2.jpg")

        # Add Thank You Slide with specific image for these templates
        add_slide(prs, "Thank You", "Thank You for Your Attention!", template_choice, image_path="Thankyou2.png")

    elif template_choice in ["Corporate6", "Creative5"]:
        # Add Q&A Slide with specific image for these templates
        add_slide(prs, "Q&A", "Questions & Answers", template_choice, image_path="QandA2.jpg")

        # Add Thank You Slide with specific image for these templates
        add_slide(prs, "Thank You", "Thank You for Your Attention!", template_choice, image_path="Thankyou3.png")

    else:
        # Add Q&A Slide with default image for other templates
        add_slide(prs, "Q&A", "Questions & Answers", template_choice, image_path="Q&As.jpg")

        # Add Thank You Slide with default image for other templates
        add_slide(prs, "Thank You", "Thank You for Your Attention!", template_choice, image_path="Thank You.jpg")



    return prs














def remove_slides(prs, indices_to_remove):
    indices_to_remove = sorted(indices_to_remove, reverse=True)
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)

    for idx in indices_to_remove:
        if idx < len(slides):
            xml_slides.remove(slides[idx])


def save_presentation(prs, filename="presentation.pptx"):
    prs.save(filename)
    print(f"Presentation saved as {filename}")


app = FastAPI()

@app.post("/generate_presentation/")
async def generate_presentation_endpoint(file: UploadFile, slide_count: int = Form(...), template_choice: str = Form(...)):
    if slide_count < 5 or slide_count > 15:
        raise HTTPException(status_code=400, detail="Slide count must be between 5 and 15.")

    try:
        if file.content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            file_content = extract_text_from_word(file.file)
        elif file.content_type == "application/pdf":
            file_content = extract_text_from_pdf(file.file)
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type.")

        # Create the presentation
        prs = create_presentation(file_content, slide_count, template_choice)

        # Remove slides 0 and 3
        remove_slides(prs, [0, 3])  # Place this after creating the presentation

        # Save the presentation
        save_presentation(prs, filename="presentation.pptx")

        # Return the presentation file for download
        return FileResponse(path="presentation.pptx", filename="presentation.pptx", media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))



if __name__ == "__main__":
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
